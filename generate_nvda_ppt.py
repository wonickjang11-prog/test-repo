"""
NVDA 재무 데이터 PPT 생성기

NVIDIA 분기별 재무 데이터를 표와 차트가 포함된 PPT로 생성합니다.

사용법:
    python generate_nvda_ppt.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np

# ── NVDA 최근 분기별 재무 데이터 (단위: 백만 달러, EPS는 달러) ──
QUARTERS = [
    "Q1 2024", "Q2 2024", "Q3 2024", "Q4 2024",
    "Q1 2025", "Q2 2025", "Q3 2025", "Q4 2025",
]
REVENUE =       [7_192, 13_507, 18_120, 22_103, 26_044, 30_040, 35_082, 39_331]
COST_OF_REV =   [2_544,  4_045,  4_720,  5_312,  5_638,  6_599,  7_606,  8_695]
GROSS_PROFIT =  [4_648,  9_462, 13_400, 16_791, 20_406, 23_441, 27_476, 30_636]
OPERATING_INC = [2_662,  6_800, 10_417, 13_615, 16_909, 19_521, 23_276, 26_033]
NET_INCOME =    [2_043,  6_188,  9_243, 12_285, 14_881, 16_599, 19_309, 22_091]
EPS =           [ 0.82,   2.48,   3.71,   4.93,   5.98,   6.67,   7.76,   8.87]
GROSS_MARGIN =  [ 64.6,   70.1,   74.0,   76.0,   78.4,   78.0,   78.3,   77.9]
OP_MARGIN =     [ 37.0,   50.3,   57.5,   61.6,   64.9,   65.0,   66.3,   66.2]
NET_MARGIN =    [ 28.4,   45.8,   51.0,   55.6,   57.1,   55.3,   55.0,   56.2]

# ── 색상 정의 ──
NVIDIA_GREEN = RGBColor(0x76, 0xB9, 0x00)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
DARK_CARD = RGBColor(0x25, 0x25, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)


def set_slide_bg(slide, color):
    """슬라이드 배경색 설정"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text,
                font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="맑은 고딕"):
    """텍스트박스를 추가합니다."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color):
    """둥근 사각형 도형을 추가합니다."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def create_chart_image(filename, chart_type, data, labels, title, colors):
    """matplotlib로 차트 이미지를 생성합니다."""
    plt.rcParams['font.family'] = 'DejaVu Sans'
    fig, ax = plt.subplots(figsize=(9, 4.5))
    fig.patch.set_facecolor('#1A1A2E')
    ax.set_facecolor('#1A1A2E')

    x = np.arange(len(labels))
    width = 0.35

    if chart_type == "bar":
        bars = ax.bar(x, data[0], width, label=data[2], color=colors[0], edgecolor='none')
        if len(data) > 3:
            bars2 = ax.bar(x + width, data[1], width, label=data[3], color=colors[1], edgecolor='none')

    elif chart_type == "line":
        for i, (d, label, color) in enumerate(zip(data[0], data[1], colors)):
            ax.plot(x, d, marker='o', linewidth=2.5, label=label, color=color, markersize=6)

    elif chart_type == "bar_line":
        ax.bar(x, data[0], width=0.5, label=data[2], color=colors[0], edgecolor='none', alpha=0.8)
        ax2 = ax.twinx()
        ax2.plot(x, data[1], marker='o', linewidth=2.5, label=data[3], color=colors[1], markersize=6)
        ax2.set_ylabel(data[3], color='white', fontsize=11)
        ax2.tick_params(colors='white', labelsize=9)
        ax2.spines['right'].set_color('#444')
        ax2.spines['left'].set_color('#444')
        ax2.spines['top'].set_visible(False)
        ax2.spines['bottom'].set_color('#444')
        ax2.legend(loc='upper left', facecolor='#25253D', edgecolor='#444',
                   labelcolor='white', fontsize=9, bbox_to_anchor=(0.75, 1.0))

    ax.set_title(title, color='white', fontsize=14, fontweight='bold', pad=15)
    ax.set_xticks(x)
    ax.set_xticklabels(labels, color='white', fontsize=9)
    ax.tick_params(colors='white', labelsize=9)

    for spine in ax.spines.values():
        spine.set_color('#444')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(axis='y', color='#333', linestyle='--', alpha=0.5)
    ax.legend(facecolor='#25253D', edgecolor='#444', labelcolor='white', fontsize=9)

    plt.tight_layout()
    plt.savefig(filename, dpi=180, facecolor='#1A1A2E', bbox_inches='tight')
    plt.close()


def slide_title_page(prs):
    """슬라이드 1: 타이틀 페이지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_BG)

    # NVIDIA 녹색 라인
    from pptx.enum.shapes import MSO_SHAPE
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(3.2), Inches(13.33), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = NVIDIA_GREEN
    line.line.fill.background()

    add_textbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
                "NVIDIA (NVDA)", font_size=48, color=NVIDIA_GREEN, bold=True)
    add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(0.8),
                "Quarterly Financial Report  |  FY2024 - FY2025",
                font_size=24, color=LIGHT_GRAY)
    add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(1.2),
                "Income Statement  /  Revenue & Profit Analysis  /  Margins & EPS Trend",
                font_size=16, color=LIGHT_GRAY)
    add_textbox(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                "Data Source: stockanalysis.com  |  Generated by Python",
                font_size=12, color=RGBColor(0x88, 0x88, 0x88))


def slide_kpi_summary(prs):
    """슬라이드 2: 핵심 KPI 요약 (최신 분기)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                "Q4 FY2025 Key Metrics", font_size=32, color=WHITE, bold=True)

    kpis = [
        ("Revenue", f"${REVENUE[-1]:,}M", f"+{(REVENUE[-1]/REVENUE[-2]-1)*100:.1f}% QoQ", NVIDIA_GREEN),
        ("Net Income", f"${NET_INCOME[-1]:,}M", f"+{(NET_INCOME[-1]/NET_INCOME[-2]-1)*100:.1f}% QoQ", ACCENT_BLUE),
        ("EPS", f"${EPS[-1]:.2f}", f"+{(EPS[-1]/EPS[-2]-1)*100:.1f}% QoQ", ACCENT_ORANGE),
        ("Gross Margin", f"{GROSS_MARGIN[-1]:.1f}%", f"vs {GROSS_MARGIN[-2]:.1f}% prev", RGBColor(0xBB, 0x86, 0xFC)),
        ("Operating Margin", f"{OP_MARGIN[-1]:.1f}%", f"vs {OP_MARGIN[-2]:.1f}% prev", RGBColor(0x03, 0xDA, 0xC6)),
        ("YoY Revenue Growth", f"+{(REVENUE[-1]/REVENUE[-5]-1)*100:.0f}%", f"vs Q4 FY2024", ACCENT_RED),
    ]

    for i, (label, value, sub, color) in enumerate(kpis):
        col = i % 3
        row = i // 3
        left = Inches(0.5 + col * 4.2)
        top = Inches(1.3 + row * 2.5)

        card = add_rounded_rect(slide, left, top, Inches(3.8), Inches(2.0), DARK_CARD)

        add_textbox(slide, left + Inches(0.3), top + Inches(0.2),
                    Inches(3.2), Inches(0.4), label,
                    font_size=14, color=LIGHT_GRAY)
        add_textbox(slide, left + Inches(0.3), top + Inches(0.7),
                    Inches(3.2), Inches(0.7), value,
                    font_size=32, color=color, bold=True)
        add_textbox(slide, left + Inches(0.3), top + Inches(1.4),
                    Inches(3.2), Inches(0.4), sub,
                    font_size=12, color=LIGHT_GRAY)


def slide_income_table(prs):
    """슬라이드 3: Income Statement 테이블"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                "Income Statement (Quarterly)", font_size=28, color=WHITE, bold=True)

    # 테이블 생성
    rows_count = 7
    cols_count = len(QUARTERS) + 1
    table_shape = slide.shapes.add_table(rows_count, cols_count,
                                          Inches(0.3), Inches(1.2),
                                          Inches(12.7), Inches(5.5))
    table = table_shape.table

    # 헤더
    header_labels = ["Metric ($ millions)"] + QUARTERS
    metrics = [
        ("Revenue", REVENUE),
        ("Cost of Revenue", COST_OF_REV),
        ("Gross Profit", GROSS_PROFIT),
        ("Operating Income", OPERATING_INC),
        ("Net Income", NET_INCOME),
        ("EPS ($)", EPS),
    ]

    def style_cell(cell, text, bg_color, font_color, bold=False, font_size=11):
        cell.text = str(text)
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color
            p.font.bold = bold
            p.font.name = "맑은 고딕"
            p.alignment = PP_ALIGN.CENTER

    # 헤더 행
    for j, label in enumerate(header_labels):
        style_cell(table.cell(0, j), label, RGBColor(0x30, 0x30, 0x50),
                   NVIDIA_GREEN, bold=True, font_size=11)

    # 데이터 행
    for i, (name, data) in enumerate(metrics):
        row_bg = DARK_CARD if i % 2 == 0 else RGBColor(0x20, 0x20, 0x38)
        style_cell(table.cell(i + 1, 0), name, row_bg, WHITE, bold=True)
        for j, val in enumerate(data):
            if name == "EPS ($)":
                txt = f"${val:.2f}"
            else:
                txt = f"${val:,}"
            style_cell(table.cell(i + 1, j + 1), txt, row_bg, WHITE)

    # 열 너비 조정
    table.columns[0].width = Inches(2.2)
    for j in range(1, cols_count):
        table.columns[j].width = Inches(1.3)


def slide_revenue_chart(prs):
    """슬라이드 4: Revenue & Net Income 차트"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                "Revenue vs Net Income", font_size=28, color=WHITE, bold=True)

    chart_file = "/tmp/nvda_revenue_chart.png"
    create_chart_image(
        chart_file, "bar",
        [REVENUE, NET_INCOME, "Revenue", "Net Income"],
        QUARTERS,
        "Revenue vs Net Income ($ millions)",
        ['#76B900', '#0096D6']
    )
    slide.shapes.add_picture(chart_file, Inches(0.5), Inches(1.0), Inches(12), Inches(6.0))
    os.remove(chart_file)


def slide_margin_chart(prs):
    """슬라이드 5: Margin 추이 차트"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                "Profitability Margins Trend", font_size=28, color=WHITE, bold=True)

    chart_file = "/tmp/nvda_margin_chart.png"
    create_chart_image(
        chart_file, "line",
        [[GROSS_MARGIN, OP_MARGIN, NET_MARGIN],
         ["Gross Margin %", "Operating Margin %", "Net Margin %"],
        ],
        QUARTERS,
        "Profitability Margins (%)",
        ['#76B900', '#FF8C00', '#0096D6']
    )
    slide.shapes.add_picture(chart_file, Inches(0.5), Inches(1.0), Inches(12), Inches(6.0))
    os.remove(chart_file)


def slide_eps_chart(prs):
    """슬라이드 6: EPS & Revenue 콤보 차트"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                "Revenue & EPS Growth", font_size=28, color=WHITE, bold=True)

    chart_file = "/tmp/nvda_eps_chart.png"
    create_chart_image(
        chart_file, "bar_line",
        [REVENUE, EPS, "Revenue ($M)", "EPS ($)"],
        QUARTERS,
        "Revenue (bars) & EPS (line)",
        ['#76B900', '#FF4545']
    )
    slide.shapes.add_picture(chart_file, Inches(0.5), Inches(1.0), Inches(12), Inches(6.0))
    os.remove(chart_file)


def slide_yoy_growth(prs):
    """슬라이드 7: YoY 성장률 비교"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                "Year-over-Year Growth", font_size=28, color=WHITE, bold=True)

    # YoY 성장률 계산 (후반 4분기 vs 전반 4분기)
    yoy_labels = QUARTERS[4:]
    rev_yoy = [(REVENUE[i+4]/REVENUE[i] - 1) * 100 for i in range(4)]
    ni_yoy = [(NET_INCOME[i+4]/NET_INCOME[i] - 1) * 100 for i in range(4)]
    eps_yoy = [(EPS[i+4]/EPS[i] - 1) * 100 for i in range(4)]

    chart_file = "/tmp/nvda_yoy_chart.png"
    create_chart_image(
        chart_file, "line",
        [[rev_yoy, ni_yoy, eps_yoy],
         ["Revenue YoY %", "Net Income YoY %", "EPS YoY %"]],
        yoy_labels,
        "Year-over-Year Growth Rate (%)",
        ['#76B900', '#0096D6', '#FF8C00']
    )
    slide.shapes.add_picture(chart_file, Inches(0.5), Inches(1.0), Inches(12), Inches(6.0))
    os.remove(chart_file)


def slide_closing(prs):
    """슬라이드 8: 마무리"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    from pptx.enum.shapes import MSO_SHAPE
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(3.4), Inches(13.33), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = NVIDIA_GREEN
    line.line.fill.background()

    add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
                "NVIDIA Financial Summary", font_size=40, color=NVIDIA_GREEN,
                bold=True, alignment=PP_ALIGN.CENTER)

    highlights = (
        f"Q4 FY2025 Revenue: ${REVENUE[-1]:,}M  |  "
        f"Net Income: ${NET_INCOME[-1]:,}M  |  "
        f"EPS: ${EPS[-1]:.2f}  |  "
        f"YoY Revenue Growth: +{(REVENUE[-1]/REVENUE[-5]-1)*100:.0f}%"
    )
    add_textbox(slide, Inches(1), Inches(4.0), Inches(11), Inches(1.0),
                highlights, font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
                "Source: stockanalysis.com  |  Generated with Python & python-pptx",
                font_size=12, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)


def generate_ppt(output_path="nvda_financial_report.pptx"):
    """NVDA 재무 보고서 PPT를 생성합니다."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("NVDA 재무 보고서 PPT 생성 중...")

    print("  [1/7] 타이틀 페이지...")
    slide_title_page(prs)

    print("  [2/7] KPI 요약...")
    slide_kpi_summary(prs)

    print("  [3/7] Income Statement 테이블...")
    slide_income_table(prs)

    print("  [4/7] Revenue vs Net Income 차트...")
    slide_revenue_chart(prs)

    print("  [5/7] Profitability Margins 차트...")
    slide_margin_chart(prs)

    print("  [6/7] Revenue & EPS 차트...")
    slide_eps_chart(prs)

    print("  [7/7] YoY Growth & Closing...")
    slide_yoy_growth(prs)
    slide_closing(prs)

    prs.save(output_path)
    print(f"\nPPT 저장 완료: {output_path}")
    print(f"총 {len(prs.slides)}장 슬라이드")
    return output_path


if __name__ == "__main__":
    generate_ppt()
